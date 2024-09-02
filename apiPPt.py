import streamlit as st
from pptx import Presentation
import requests
import os

# Set your Google API key here
API_KEY = key

# Function to generate text based on a given topic using a hypothetical Google API
def generate_text_from_topic(topic):
    url = f"https://language.googleapis.com/v1/documents:analyzeEntities?key={API_KEY}"
    headers = {
        "Content-Type": "application/json"
    }
    data = {
        "document": {
            "type": "PLAIN_TEXT",
            "content": f"Write a detailed summary or introduction about the topic: {topic}"
        },
        "encodingType": "UTF8"
    }
    
    response = requests.post(url, headers=headers, json=data)
    result = response.json()
    
    # This is just an example. You'd parse the response based on the actual API response format.
    generated_text = "Generated text based on your topic."  # Placeholder
    return generated_text

# Function to create a PowerPoint presentation from text
def create_ppt_from_text(slides_data):
    prs = Presentation()

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using layout 1 (Title and Content)
        title, content = slide_data
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

    return prs

# Streamlit app
st.title("Topic to PPT Generator")

# Step 1: Input the topic name
st.header("Step 1: Input a Topic")
topic = st.text_input("Enter the topic")

# Step 2: Generate text automatically from the topic
st.header("Step 2: Generate Text and Create PPT")
if st.button("Generate PPT"):
    if topic:
        # Generate text based on the topic
        generated_text = generate_text_from_topic(topic)
        
        # Split the generated text into multiple slides (for simplicity, break text into chunks of 100 words per slide)
        slide_texts = generated_text.split(". ")  # Split text by sentences for better distribution
        slides_data = []
        for i, chunk in enumerate(slide_texts):
            slides_data.append((f"Slide {i+1}: {topic}", chunk))

        # Create the PowerPoint presentation from the text
        presentation = create_ppt_from_text(slides_data)
        ppt_filename = "generated_presentation.pptx"
        presentation.save(ppt_filename)

        # Allow the user to download the PPT
        with open(ppt_filename, "rb") as ppt_file:
            st.download_button(
                label="Download PPT",
                data=ppt_file,
                file_name=ppt_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        
        os.remove(ppt_filename)
    else:
        st.warning("Please enter a topic to generate the PPT.")

# Step 3: Download PPT
st.header("Step 3: Download the Generated PPT")
st.write("After generating the PPT, you can download it using the button above.")
